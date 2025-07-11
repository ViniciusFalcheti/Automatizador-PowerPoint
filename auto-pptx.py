from pptx import Presentation
from pptx.util import Pt
from docx import Document
import tkinter as tk
from tkinter import simpledialog
from tkinter import messagebox
import re
import traceback

class PptxAutomationFromDocx:

    def __init__(self, arquivo_docx):
        try:
            self.prs = Presentation('PADRAO-CULTO-ONLINE.pptx')  # Utilize seu modelo real aqui
            self.dados = self.ler_arquivo_word(arquivo_docx)
            self.culto = 'Yes'  # Você pode melhorar isso para ser lido do Word também se quiser
            self.tema = self.escolher_tema()

            self.montar_apresentacao()
        except Exception as e:
            if '[Errno 13] Permission denied' in str(e):
                self.exibir_mensagem_erro("O arquivo está aberto, feche-o e tente novamente.")
            else:
                self.exibir_mensagem_erro(e)

            traceback.print_exc()

    def escolher_tema(self):
        root = tk.Tk()
        root.withdraw()
        tema = simpledialog.askinteger("Escolher Tema", "Escolha o tema:\n1 - Padrão Online\n2 - Padrão manhã\n3 - Padrão Yes:", minvalue=1, maxvalue=3)
        
        if tema is None:
            tema = 0
        elif tema == 1:
            tema = 0
        elif tema == 2:
            tema = 4
        elif tema == 3:
            tema = 2

        return tema
    
    def exibir_mensagem_sucesso(self):
        messagebox.showinfo("Sucesso", "A apresentação PPTX foi gerada com sucesso!")

    def exibir_mensagem_erro(self, error=None):
        if error:
            messagebox.showerror("Erro", f"Ocorreu um erro: {error}", icon='error')
        else:
            messagebox.showinfo("Erro", "Não foi possível gerar o arquivo, verifique se ele está aberto.", icon='error')

    def ler_arquivo_word(self, caminho):
        doc = Document(caminho)
        dados = {
            'titulo': '',
            'pregador': '',
            'versiculo_chave': None,
            'pontos': [],
            'frase': None
        }

        ponto_atual = None
        versiculo_chave_pego = False
        texto_chave = ""  # Variável para armazenar o texto chave
        el_anterior = None  # Variável para armazenar o elemento anterior
        regex = re.compile(r'^((\d{1,2})|([\u00B2\u00B3\u00B9\u2070-\u2079]{1,2}))\s?') # Regex para identificar textos que começam com números ou sobrescritos

        for para in doc.paragraphs:
            texto = para.text.strip()
            if not texto:
                continue

            if texto.lower().startswith(('título:', 'titulo:')):
                dados['titulo'] = texto.replace('Título:', '').replace('Titulo:', '').strip()
                el_anterior = 'titulo'

            elif texto.lower().startswith(('pregador:', 'pregadora:', 'pastor:', 'pastora:')):
                dados['pregador'] = texto.replace('Pregador:', '').strip()
                el_anterior = 'pregador'

            elif texto.lower().startswith(('versículo chave:', 'versiculo chave:')):
                referencia = texto.replace('Versículo chave:', '').replace('Versiculo chave:', '').strip()
                dados['versiculo_chave'] = {'referencia': referencia, 'texto': []}
                versiculo_chave_pego = True
                el_anterior = 'versiculo_chave'

            # Coleta o texto chave (o conteúdo que aparece após "Texto chave:")
            elif texto.lower().startswith('texto chave:') and versiculo_chave_pego:
                texto_chave = texto.replace('Texto chave:', '').strip()
                dados['versiculo_chave']['texto'].append(texto_chave)
                el_anterior = 'texto_chave'

            # Quando encontramos um ponto, começamos a adicionar seus versículos
            elif texto.lower().startswith('ponto'):
                if ponto_atual:
                    dados['pontos'].append(ponto_atual)
                ponto_atual = {'texto': texto.split(':', 1)[1].strip(), 'versiculos': [], 'subtitulo': '', 'frases': []}
                el_anterior = 'ponto'

            elif texto.lower().startswith(('subtítulo:', 'subtitulo:')) and ponto_atual:
                ponto_atual['subtitulo'] = texto.replace('Subtítulo:', '').replace('Subtitulo:', '').strip()
                el_anterior = 'subtitulo'

            elif texto.lower().startswith(('versículo:', 'versiculo:')):
                if ponto_atual is not None:
                    versiculo = {'referencia': texto.replace('Versículo:', '').replace('Versiculo:', '').strip(), 'texto': []}
                    ponto_atual['versiculos'].append(versiculo)
                    el_anterior = 'versiculo'

            # Adiciona o texto de cada versículo ao ponto correto, mas NÃO mistura com o versículo chave
            elif texto.lower().startswith('texto:'):
                if ponto_atual and ponto_atual['versiculos']:
                    ponto_atual['versiculos'][-1]['texto'].append(texto.replace('Texto:', '').strip())
                    el_anterior = 'texto'

            elif texto.lower().startswith('frase:'):
                frase = texto.replace('Frase:', '').replace('frase:', '').strip()
                ponto_atual['frases'].append(frase)
                el_anterior = 'frase'

            # elif texto.lower().startswith('frase:'):
            #     dados['frase'] = texto.replace('Frase:', '').strip()
            #     el_anterior = 'frase'
                

            elif regex.match(texto):

                if el_anterior == 'texto_chave' and dados['versiculo_chave']:
                    dados['versiculo_chave']['texto'].append(texto)

                elif el_anterior == 'texto' and ponto_atual and ponto_atual['versiculos']:
                    ponto_atual['versiculos'][-1]['texto'].append(texto)

        if ponto_atual:
            dados['pontos'].append(ponto_atual)

        return dados

    def montar_apresentacao(self):
        self.criar_slide_titulo()

        if self.dados['versiculo_chave']:
            self.criar_slides_de_versiculos(self.dados['versiculo_chave']['referencia'], self.dados['versiculo_chave']['texto'])

        for i, ponto in enumerate(self.dados['pontos'], start=1):
            self.criar_slide_ponto(ponto['texto'], i, ponto['subtitulo']) #ponto['subtitulo']
            for versiculo in ponto['versiculos']:
                self.criar_slides_de_versiculos(versiculo['referencia'], versiculo['texto'])
            for frase in ponto['frases']:
                self.criar_slide_frase(frase)
        
        # if self.tema == 4 and self.dados['frase']: # Tema 4 - Padrão Manhã
        #     self.criar_slide_frase(self.dados['frase'])

        nome_arquivo = f"{self.dados['titulo'].strip()} - {self.dados['pregador'].strip()}.pptx"
        nome_arquivo = self.limpar_nome_arquivo(nome_arquivo)

        self.prs.save(nome_arquivo)
        print(f"Apresentação salva como: {nome_arquivo}")
        self.exibir_mensagem_sucesso()

    def criar_slide_titulo(self):
        slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[0])
        titulo = slide.placeholders[10] if self.tema == 0 else slide.shapes.title
        # titulo = slide.shapes.title
        titulo.text = self.dados['titulo'].strip().upper()

        if self.tema == 4: # Tema 4 - Padrão Manhã
            textTextoChave = slide.placeholders[11]
            textTextoChave.text = self.dados['versiculo_chave']['referencia'].strip()

            textPregador = slide.placeholders[10]
            textPregador.text = self.dados['pregador'].strip()

        if self.tema != 4:
            self.change_pregador_name_to_bold(self.dados['pregador'].strip())

        tituloParametro = self.dados['titulo'].strip() if self.tema == 4 else f"{self.dados['titulo'].strip()} {self.dados['pregador'].strip()}"

        titulo.text_frame.paragraphs[0].font.size = self.ajustar_tamanho_fonte_por_texto(tituloParametro, tipo='titulo')

    def criar_slide_ponto(self, texto_ponto, numero, subtitulo_Ponto=None):
        if self.tema == 0: # Tema 0 - Padrão Online
            slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[numero+1])
            textPonto = slide.shapes.title
            textPonto.text = texto_ponto.strip()
        elif self.tema == 4: # Tema 4 - Padrão Manhã
            if subtitulo_Ponto:
                slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[3])
            else:
                slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[2])
            textPonto = slide.shapes.title
            textPonto.text = f"{numero}. {texto_ponto.upper().strip()}"
        else:
            slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[2])
            textPonto = slide.shapes.title
            textPonto.text = f"{numero}. {texto_ponto.strip()}"

        textPonto.text_frame.paragraphs[0].font.size = self.ajustar_tamanho_fonte_por_texto(f"{numero}. {texto_ponto.strip()}", tipo='ponto')

        if self.tema == 4: # Tema 4 - Padrão Manhã
            textTituloPalavra = slide.placeholders[11]
            textTituloPalavra.text = self.dados['titulo'].strip().upper()
        
            if subtitulo_Ponto:
                textSubtitulo = slide.placeholders[10]
                textSubtitulo.text = subtitulo_Ponto.strip()

    def criar_slide_frase(self, frase):
        if self.tema == 4 : # Tema 4 - Padrão Manhã
            slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[4])
            frasePlaceHolder = slide.shapes.title
            frasePlaceHolder.text = frase.upper().strip()
        elif self.tema == 0: # Tema 0 - Padrão Online
            slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[12])
            frasePlaceHolder = slide.placeholders[11]
            frasePlaceHolder.text = frase.strip()

        frasePlaceHolder.text_frame.paragraphs[0].font.size = self.ajustar_tamanho_fonte_por_texto(frase, tipo='frase')

    def criar_slides_de_versiculos(self, referencia, texto):
        versiculos = self.agrupar_versiculos_por_paragrafo(texto)

        if self.tema == 4: # Tema 4 - Padrão Manhã
            agrupados = [versiculos[i:i+3] for i in range(0, len(versiculos), 3)]
        else:
            agrupados = [[v] for v in versiculos]
        
        for grupo in agrupados:
            texto_slide = "".join(grupo)
            self.criar_slide_versiculo(referencia, texto_slide)

    def criar_slide_versiculo(self, referencia, texto):
        slide = self.prs.slides.add_slide(self.prs.slide_masters[self.tema].slide_layouts[1])
        if self.tema == 0:
            textVer = slide.placeholders[11]
            textVer.text = texto

            textRef = slide.placeholders[10]
            textRef.text = referencia
        elif self.tema == 2:
            textRef = slide.shapes.title
            textRef.text = referencia

            textVer = slide.placeholders[1]
            textVer.text = texto
        else:
            textVer = slide.placeholders[11]
            textVer.text = texto

            textRef = slide.placeholders[12]
            textRef.text = referencia

        textRef.text_frame.paragraphs[0].font.size = self.ajustar_tamanho_fonte_por_texto(referencia, tipo='refVersiculo')
        textVer.text_frame.paragraphs[0].font.size = self.ajustar_tamanho_fonte_por_texto(texto, tipo='versiculo')

    def eh_novo_versiculo(self, linha):
        return re.match(r'^(\d{1,2}|[\u00B2\u00B3\u00B9\u2070-\u2079]{1,2})\s', linha.strip()) is not None

    def agrupar_versiculos_por_paragrafo(self, texto):
        """
        Agrupa versículos do ponto e texto chave de forma separada, sem misturar
        com base no padrão de numeração ou sobrescrito.
        """
        versiculos = []
        
        if isinstance(texto, list):
            linhas = texto  # Caso o texto já seja uma lista de linhas
        else:
            linhas = texto.split('\n')  # Caso o texto venha como uma string

        atual = ''
        for linha in linhas:
            # Verifica se a linha inicia com um versículo (numérico ou sobrescrito)
            if self.eh_novo_versiculo(linha):
                if atual:
                    versiculos.append(atual.strip())  # Se 'atual' tem texto, adiciona à lista de versículos
                atual = linha.strip()  # Inicia um novo versículo com a linha atual
            else:
                atual += " " + linha.strip()  # Adiciona mais texto ao versículo atual
        
        if atual:
            versiculos.append(atual.strip())  # Adiciona o último versículo

        return versiculos

    def change_pregador_name_to_bold(self, pregador):
        slide = self.prs.slides[(0)]
        # title = slide.shapes.title
        title = slide.placeholders[10] if self.tema == 0 else slide.shapes.title
        tf = title.text_frame

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f' {pregador}'
        font = run.font
        font.bold = True

    def ajustar_tamanho_fonte_por_texto(self, texto, tipo='titulo'):
        comprimento = len(texto)

        if tipo == 'titulo':
            if self.tema == 0: # Tema 0 - Padrão Online
                if comprimento <= 30:
                    return Pt(32)
                elif comprimento <= 85:
                    return Pt(24)
                else:
                    return Pt(18)
            if self.tema == 2: # Tema 2 - Padrão Yes
                if comprimento <= 25:
                    return Pt(40)
                if comprimento <= 35:
                    return Pt(32)
                if comprimento <= 45:
                    return Pt(30)
                else:
                    return Pt(24)
            elif self.tema == 4: # Tema 4 - Padrão Manhã
                if comprimento <= 20:
                    return Pt(100)
                elif comprimento <= 30:
                    return Pt(80)
                elif comprimento <= 50:
                    return Pt(60)
                else:
                    return Pt(48)
        elif tipo == 'versiculo':
            if self.tema == 0: # Tema 0 - Padrão Online
                if comprimento <= 190:
                    return Pt(24)
                elif comprimento <= 310:
                    return Pt(20)
                elif comprimento <= 420:
                    return Pt(18)
                elif comprimento <= 580:
                    return Pt(16)
                else:
                    return Pt(14)
            elif self.tema == 2: # Tema 2 - Padrão Yes
                if comprimento <= 300:
                    return Pt(18)
                elif comprimento <= 580:
                    return Pt(16)
                else:
                    return Pt(14)
            elif self.tema == 4: # Tema 4 - Padrão Manhã
                if comprimento <= 380:
                    return Pt(30)
                elif comprimento <= 460:
                    return Pt(28)
                else:
                    return Pt(24)
        elif tipo == 'ponto':
            if self.tema == 0: # Tema 0 - Padrão Online
                if comprimento <= 30:
                    return Pt(32)
                elif comprimento <= 85:
                    return Pt(24)
                else:
                    return Pt(18)
            if self.tema == 2: # Tema 2 - Padrão Yes
                if comprimento <= 40:
                    return Pt(36)
                elif comprimento <= 85:
                    return Pt(26)
                else:
                    return Pt(22)
            elif self.tema == 4: # Tema 4 - Padrão Manhã
                if comprimento <= 25:
                    return Pt(70)
                elif comprimento <= 40:
                    return Pt(54)
                else:
                    return Pt(48)
        elif tipo == 'refVersiculo':
            if self.tema == 0: # Tema 0 - Padrão Online
                if comprimento <= 13:
                    return Pt(24)
                elif comprimento <= 20:
                    return Pt(20)
                else:
                    return Pt(18)
            if self.tema == 2: # Tema 2 - Padrão Yes
                if comprimento <= 15:
                    return Pt(18)
                elif comprimento <= 20:
                    return Pt(16)
                else:
                    return Pt(14)
            if self.tema == 4: # Tema 4 - Padrão Manhã
                if comprimento <= 18:
                    return Pt(26)
                elif comprimento <= 25:
                    return Pt(20)
                else:
                    return Pt(18)
        elif tipo == 'frase':
            if self.tema == 0: # Tema 0 - Padrão Online
                if comprimento <= 80:
                    return Pt(28)
                elif comprimento <= 125:
                    return Pt(24)
                else:
                    return Pt(20)
            elif self.tema == 4: # Tema 4 - Padrão Manhã
                if comprimento <= 20:
                    return Pt(100)
                elif comprimento <= 35:
                    return Pt(80)
                elif comprimento <= 55:
                    return Pt(72)
                elif comprimento <= 90:
                    return Pt(64)
                elif comprimento <= 130:
                    return Pt(48)
                else:
                    return Pt(40)
        
    def limpar_nome_arquivo(self, nome):
        # Remove caracteres inválidos
        nome_limpo = re.sub(r'[<>:"/\\|?*]', '', nome)

        # Remove espaços ou pontos no final
        nome_limpo = nome_limpo.rstrip(' .')

        # Substitui por 'arquivo' se o nome ficar vazio ou só com espaços
        if not nome_limpo.strip():
            return 'arquivo'

        return nome_limpo
    
if __name__ == "__main__":
    caminho_docx = 'entrada.docx'  # Nome do arquivo Word que você criou
    PptxAutomationFromDocx(caminho_docx)